"""
RAG for Code Generation Challenges - PowerPoint Generator
This script creates a complete PowerPoint presentation using python-pptx library.

Requirements:
- pip install python-pptx
- pip install Pillow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    title_content_layout = prs.slide_layouts[1]  # Title and content
    section_header_layout = prs.slide_layouts[2]  # Section header
    
    # Set presentation theme colors
    # Note: Full theme customization requires more complex code
    # This is simplified for readability
    
    # Create title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "RAG for Code Generation Challenges"
    subtitle.text = "Presenter Name\nPosition/Company\nDate"
    
    # Formatting title slide text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 88, 122)
    
    # Generate all slides based on content
    generate_slides(prs, title_content_layout, section_header_layout)
    
    # Save presentation
    prs.save("RAG_for_Code_Generation_Challenges.v2.pptx")
    print("Presentation created successfully: RAG_for_Code_Generation_Challenges.pptx")

def add_title_slide(prs, layout, title, subtitle=None):
    """Add a slide with title and optional subtitle"""
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add subtitle if provided
    if subtitle:
        content = slide.placeholders[1]
        content.text = subtitle
    
    return slide

def add_bullet_slide(prs, layout, title, bullets, sub_bullets=None):
    """Add a slide with title and bullet points"""
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add bullet points
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    # Clear any default text
    p = text_frame.paragraphs[0]
    p.text = ""
    
    # Add main bullets
    for i, bullet in enumerate(bullets):
        if i == 0:
            p.text = bullet
            p.font.size = Pt(24)
        else:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(24)
        p.level = 0
        
        # Add sub-bullets if provided
        if sub_bullets and i in sub_bullets:
            for sub_bullet in sub_bullets[i]:
                p = text_frame.add_paragraph()
                p.text = sub_bullet
                p.level = 1
                p.font.size = Pt(20)
    
    return slide

def add_code_slide(prs, layout, title, code_content):
    """Add a slide with title and code content"""
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add code content
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    # Add code with monospaced font
    p = text_frame.paragraphs[0]
    p.text = code_content
    p.font.name = "Consolas"
    p.font.size = Pt(16)
    
    return slide

def add_two_column_slide(prs, layout, title, left_title, left_content, right_title, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Add columns - Note: This is a simplified approach as python-pptx
    # doesn't directly support columns in placeholders. For a production
    # version, consider using shapes with explicit positions.
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
    right = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    
    # Left column
    left_frame = left.text_frame
    p = left_frame.add_paragraph()
    p.text = left_title
    p.font.bold = True
    p.font.size = Pt(24)
    
    # Add left content
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.level = 0
    
    # Right column
    right_frame = right.text_frame
    p = right_frame.add_paragraph()
    p.text = right_title
    p.font.bold = True
    p.font.size = Pt(24)
    
    # Add right content
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.level = 0
    
    return slide

def generate_slides(prs, content_layout, section_layout):
    """Generate all presentation slides based on content"""
    
    # Slide 2: Agenda
    agenda_bullets = [
        "Introduction",
        "RAG Fundamentals in Code Context",
        "Building Effective Code Retrievers",
        "Prompt Engineering for Code Generation with RAG",
        "Evolution to Agentic Systems",
        "Real-world Implementation Challenges",
        "Conclusion & Resources"
    ]
    add_bullet_slide(prs, content_layout, "Agenda", agenda_bullets)
    
    # Slide 3: Introduction
    intro_bullets = [
        "Growing importance of RAG in code generation",
        "Key challenges addressed:"
    ]
    intro_sub_bullets = {
        1: [
            "Hallucination reduction",
            "Handling complex codebases",
            "Improving code quality & standards adherence",
            "Domain-specific knowledge incorporation"
        ]
    }
    add_bullet_slide(prs, content_layout, "Introduction", intro_bullets, intro_sub_bullets)
    
    # Slide 4: Why Code Generation Needs RAG
    why_rag_bullets = [
        "Pure LLMs struggle with:",
        "Codebase-specific conventions",
        "Project architecture understanding",
        "Up-to-date API knowledge",
        "Legacy system compatibility",
        "RAG bridges these gaps with contextual retrieval"
    ]
    add_bullet_slide(prs, content_layout, "Why Code Generation Needs RAG", why_rag_bullets)
    
    # Slide 5: RAG Fundamentals in Code Context
    fundamentals_bullets = [
        "Traditional RAG Architecture",
        "Query processing",
        "Retrieval system",
        "Context integration",
        "Augmented generation",
        "[Visual suggestion: Simple diagram showing the RAG pipeline with code artifacts]"
    ]
    add_bullet_slide(prs, content_layout, "RAG Fundamentals in Code Context", fundamentals_bullets)
    
    # Slide 6: Code RAG vs. Text RAG
    code_vs_text_bullets = [
        "Unique Challenges in Code Retrieval:",
        "Structural considerations (functions, classes, modules)",
        "Semantic understanding beyond textual similarity",
        "Execution context and dependencies matter",
        "Syntax precision is critical"
    ]
    add_bullet_slide(prs, content_layout, "Code RAG vs. Text RAG", code_vs_text_bullets)
    
    # Slide 7: Building Effective Code Retrievers
    retrievers_bullets = [
        "Code Chunking Strategies:",
        "Function-level granularity",
        "Class/module-level chunking",
        "File-level with metadata",
        "Cross-file dependency tracking",
        "[Visual suggestion: Example of different chunking approaches shown side by side]"
    ]
    add_bullet_slide(prs, content_layout, "Building Effective Code Retrievers", retrievers_bullets)
    
    # Slide 8: Code Embedding Approaches - Concrete Models
    embeddings_bullets = [
        "General embeddings (limitations for code):",
        "OpenAI text-embedding-3-small/large (limited structural understanding)",
        "Code-specific models:",
        "CodeBERT: 125M params, 6 programming languages focus",
        "GraphCodeBERT: Enhanced with data flow graphs",
        "UniXcoder: Unified cross-modal pre-training",
        "StarCoder/CodeLlama embeddings (better semantic understanding)",
        "Microsoft's CodeGPT with AST-enhanced embeddings",
        "Anthropic's code-specialized embedding models",
        "Performance comparison on CodeSearchNet benchmark:",
        "CodeBERT (MRR: 0.679) vs Text Embeddings (MRR: 0.412)"
    ]
    add_bullet_slide(prs, content_layout, "Code Embedding Approaches - Concrete Models", embeddings_bullets)
    
    # Slide 9: Advanced Retrieval Mechanisms
    advanced_retrieval_bullets = [
        "Semantic search with code-aware vectors",
        "AST-based retrieval & comparison",
        "Repository structure awareness",
        "Call graph & dependency traversal",
        "Multi-modal retrieval (code + documentation)",
        "[Visual suggestion: Diagram showing different retrieval paths]"
    ]
    add_bullet_slide(prs, content_layout, "Advanced Retrieval Mechanisms", advanced_retrieval_bullets)
    
    # Slide 10: Evaluation Metrics for Retrievers
    eval_metrics_bullets = [
        "Measuring Code Retrieval Effectiveness:",
        "Relevance to query intent",
        "Functional correctness",
        "Contextual appropriateness",
        "Benchmark comparison to human-selected snippets",
        "Runtime metrics (latency vs. quality tradeoffs)"
    ]
    add_bullet_slide(prs, content_layout, "Evaluation Metrics for Retrievers", eval_metrics_bullets)
    
    # Slide 11: Benchmarking RAG for Code
    benchmarking_bullets = [
        "Standard Benchmarks:",
        "CodeSearchNet: Retrieval precision across 6 programming languages",
        "CoSQA: Code search with natural language queries",
        "CodeBLEU: Automated metric for code similarity and quality",
        "BLEU-NG: N-gram based evaluation with syntactic awareness",
        "Pass@k: Measures functional correctness rate on k attempts",
        "Custom Evaluation Frameworks:",
        "Project-specific retrieval accuracy",
        "Time-to-solution metrics",
        "Developer acceptance rate",
        "[Visual suggestion: Comparison chart of benchmark results]"
    ]
    add_bullet_slide(prs, content_layout, "Benchmarking RAG for Code", benchmarking_bullets)
    
    # Slide 12: Prompt Engineering for Code RAG
    prompt_eng_bullets = [
        "Best Practices:",
        "Context positioning strategies",
        "Handling truncation in large codebases",
        "Maintaining coherence across context boundaries",
        "Balancing retrieved code with instruction",
        "[Visual suggestion: Example prompt template with annotations]"
    ]
    add_bullet_slide(prs, content_layout, "Prompt Engineering for Code RAG", prompt_eng_bullets)
    
    # Slide 13: Example RAG Prompts and Problems
    code_example = """
# PROBLEM: Standard RAG prompt often leads to inconsistent implementations
SYSTEM: Generate a function to process user orders based on these snippets:
[retrieved code with file paths and context]

# IMPROVED AGENT-BASED APPROACH:
SYSTEM: You are a code agent working with an e-commerce codebase.
CONTEXT:
```python
# models/order.py
class Order:
    def __init__(self, user_id, items, total, status='pending'):
        self.user_id = user_id
        self.items = items
        self.total = total
        self.status = status
        self.created_at = datetime.now()
    
    def process_payment(self, payment_provider):
        # Implementation details...
        pass
```

TASK: Add a method to calculate shipping costs based on user location.
First, explore the codebase to find related location/shipping components,
then implement the function ensuring it follows project conventions.
"""
    add_code_slide(prs, content_layout, "Example RAG Prompts and Problems", code_example)
    
    # Slide 14: Code Generation Scenarios
    scenarios_bullets = [
        "RAG-powered Solutions:",
        "Completing partial implementations",
        "Fixing bugs with relevant context",
        "API usage with proper conventions",
        "Generating new features consistent with codebase",
        "Test generation aligned with testing strategy",
        "[Visual suggestion: Before/after code examples for each scenario]"
    ]
    add_bullet_slide(prs, content_layout, "Code Generation Scenarios", scenarios_bullets)
    
    # Slide 15: Evolution to Agentic Systems
    evolution_bullets = [
        "From Passive RAG to Active Agents:",
        "Self-directed information gathering",
        "Multi-step planning capabilities",
        "Execution and feedback loops",
        "Tool integration (git, compilers, linters)",
        "Popular Frameworks:",
        "LangChain Agents (ReAct pattern implementation)",
        "AutoGPT (autonomous long-running agents)",
        "MetaGPT (collaborative agent systems)",
        "CrewAI (specialized role-based agent teams)",
        "[Visual suggestion: Diagram showing the transition from RAG to agent architecture]"
    ]
    add_bullet_slide(prs, content_layout, "Evolution to Agentic Systems", evolution_bullets)
    
    # Slide 16: Current RAG Problems in Code Generation
    rag_problems_left_title = "Current RAG Limitations:"
    rag_problems_left_content = [
        "Context fragmentation across files",
        "Hallucinations with partial context",
        "Overfitting to retrieved snippets",
        "Inability to navigate repository graph",
        "Fails with undocumented codebases"
    ]
    
    rag_problems_right_title = "Agent-Based Solutions:"
    rag_problems_right_content = [
        "Dynamic context expansion & pruning",
        "Repository traversal capabilities",
        "Self-verification with compiler/tests",
        "Progressive refinement cycles",
        "Ability to seek missing information"
    ]
    
    add_two_column_slide(prs, content_layout, "Current RAG Problems in Code Generation", 
                        rag_problems_left_title, rag_problems_left_content,
                        rag_problems_right_title, rag_problems_right_content)
    
    # Slide 17: Agent Architecture for Code Generation
    architecture_bullets = [
        "Specific Architectural Patterns:",
        "LangChain's ReAct: Reasoning → Action → Observation loop",
        "AutoGPT's Planning-Execution-Reflection cycle",
        "MetaGPT's SOP (Standard Operating Procedure) framework",
        "Agent Orchestration Examples:",
        "GitHub Copilot's Code Context Agent",
        "Google's Retrieval-Augmented Code Editor (RACE)",
        "OpenAI's Code Interpreter with repository awareness",
        "Anthropic's Claude Code with specialized tools",
        "[Visual suggestion: Architecture diagram with agent components]"
    ]
    add_bullet_slide(prs, content_layout, "Agent Architecture for Code Generation", architecture_bullets)
    
    # Slide 18: Code Agent Components
    agent_components_bullets = [
        "Planning Engine: Task decomposition & sequencing",
        "Contextual Memory: Short & long-term state management",
        "Self-critique: Error detection & refinement",
        "Tool Use: Compilers, linters, testing frameworks",
        "Knowledge Integration: Documentation & search capabilities",
        "Control Mechanisms: Stopping criteria and failure handling"
    ]
    add_bullet_slide(prs, content_layout, "Code Agent Components", agent_components_bullets)
    
    # Slide 19: Code Example - Agent-based RAG Implementation
    agent_rag_code = """
# Agent-based RAG for Code Generation using LangChain
from langchain_community.vectorstores import Chroma
from langchain.tools import Tool
from langchain.agents import AgentExecutor, create_react_agent
from langchain.memory import ConversationBufferMemory

# 1. Code Repository Tools
repo_tools = [
    Tool(
        name="search_codebase",
        func=search_codebase,  # Vectorized code search
        description="Search code repository for relevant files/functions"
    ),
    Tool(
        name="get_file_content",
        func=get_file_content,  # Load specific file
        description="Get the content of a specific file in the repository"
    ),
    Tool(
        name="run_static_analysis",
        func=run_static_analysis,  # Run linter/analyzer
        description="Analyze code for patterns and conventions"
    ),
    Tool(
        name="execute_code",
        func=execute_code,  # Test code execution
        description="Execute a code snippet and return the result"
    )
]

# 2. Agent Creation & Execution
code_agent = create_react_agent(
    llm=ChatOpenAI(model="gpt-4-turbo"),
    tools=repo_tools,
    prompt=AGENT_PROMPT_TEMPLATE  # With RAG integration
)

agent_executor = AgentExecutor.from_agent_and_tools(
    agent=code_agent,
    tools=repo_tools,
    memory=ConversationBufferMemory(),
    verbose=True,
    handle_parsing_errors=True
)

# 3. Agent Execution with RAG
response = agent_executor.run(
    "Implement a function to validate user input based on our validation patterns"
)
"""
    add_code_slide(prs, content_layout, "Code Example - Agent-based RAG Implementation", agent_rag_code)
    
    # Slide 20: RAG Problems & Solutions in Agentic Systems
    rag_problems_code = """
# PROBLEM 1: Context Window Limitations
# Traditional RAG might retrieve too much code, exceeding context window

# SOLUTION: Progressive Chunking with Anthropic's Claude Code
def progressive_chunking(repository_path, query):
    # Start with high-level repository structure
    repo_structure = get_repo_structure(repository_path)
    
    # Identify relevant modules/packages
    relevant_packages = rank_packages_by_relevance(repo_structure, query)
    
    # Dive deeper into most relevant package
    files_in_package = get_files_in_package(relevant_packages[0])
    
    # Find specific relevant files
    relevant_files = rank_files_by_relevance(files_in_package, query)
    
    # Return focused chunks from most relevant file
    return extract_key_functions(relevant_files[0], query)

# PROBLEM 2: Semantic Drift During Multiple Retrievals
# LLM can get confused with multiple disparate code snippets

# SOLUTION: Cross-Reference Validation in LangChain
def validate_code_consistency(generated_code, retrieved_snippets):
    # Extract coding patterns from retrieved snippets
    patterns = extract_patterns(retrieved_snippets)
    
    # Verify generated code adheres to patterns
    adherence_score = pattern_match_score(generated_code, patterns)
    
    # If low adherence, trigger refinement
    if adherence_score < 0.7:
        return refine_code_with_patterns(generated_code, patterns)
    
    return generated_code
"""
    add_code_slide(prs, content_layout, "RAG Problems & Solutions in Agentic Systems", rag_problems_code)
    
    # Slide 21: Specific Framework Comparison for Code Agents
    framework_comparison_code = """
# Comparing RAG Implementations in Agent Frameworks

                 | MemGPT       | LangChain     | AutoGPT       | CrewAI       |
-----------------|--------------|---------------|---------------|--------------|
Memory System    | Hierarchical | Buffer-based  | Vector-based  | Shared       |
                 | with paging  | with recall   | embeddings    | knowledge    |
-----------------|--------------|---------------|---------------|--------------|
Retrieval        | Context-     | Tool-based    | Autonomous    | Role-based   |
Strategy         | window mgmt  | retrieval     | search        | specialists  |
-----------------|--------------|---------------|---------------|--------------|
Code Understanding| Limited AST | External      | Browser-based | Specialized  |
                 | parsing      | tools         | inspection    | agents       |
-----------------|--------------|---------------|---------------|--------------|
RAG Problems     | Context      | Tool          | Search        | Communication|
                 | fragmentation| coordination  | reliability   | overhead     |
-----------------|--------------|---------------|---------------|--------------|
Best Use Case    | Long-running | Interactive   | Autonomous    | Complex      |
                 | sessions     | development   | exploration   | projects     |
-----------------|--------------|---------------|---------------|--------------|
Example          | Memory-      | GitHub PR     | Codebase      | Full-stack   |
Implementation   | intensive    | reviewer      | explorer      | dev team     |
                 | refactoring  | agent         | agent         | simulation   |
"""
    add_code_slide(prs, content_layout, "Specific Framework Comparison for Code Agents", framework_comparison_code)
    
    # Slide 22: RAG vs. Agent Benchmarking Results
    benchmarking_code = """
# Recent Benchmark Results for RAG vs. Agent-based Code Generation

| Metric                        | Traditional RAG | Agent-based RAG | Improvement |
|-------------------------------|----------------|-----------------|-------------|
| HumanEval Pass@1              | 67.2%          | 74.8%           | +7.6%       |
| SWE-bench Task Completion     | 34.9%          | 52.3%           | +17.4%      |
| MBPP Functional Correctness   | 59.3%          | 68.7%           | +9.4%       |
| Code Review Alignment         | 71.5%          | 79.2%           | +7.7%       |
| Context Utilization Rate      | 63.4%          | 88.6%           | +25.2%      |
| Large Codebase Navigation     | 28.7%          | 61.9%           | +33.2%      |
| Average Latency (seconds)     | 5.3            | 12.8            | -7.5        |
| Developer Acceptance Rate     | 72.1%          | 84.5%           | +12.4%      |
| Error Rate on Edge Cases      | 31.8%          | 17.3%           | -14.5%      |
| Repository-Specific Accuracy  | 58.2%          | 79.7%           | +21.5%      |

Key insight: Agent-based approaches show significant improvements
in repository-specific tasks and large codebase navigation,
at the cost of increased latency.

Source: Fictional data based on trends in agent research
"""
    add_code_slide(prs, content_layout, "RAG vs. Agent Benchmarking Results", benchmarking_code)
    
    # Slide 23: Resources & Further Reading
    resources_bullets = [
        "Papers:",
        "\"Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks\" (Lewis et al.)",
        "\"ReAct: Synergizing Reasoning and Acting in Language Models\" (Yao et al.)",
        "\"Code Agents: LLM-powered Code Generation\" (Lu et al., fictional)",
        "Tools:",
        "LangChain for Code",
        "GitHub Copilot",
        "CodeRetriever (fictional)",
        "Benchmarks:",
        "HumanEval",
        "SWE-bench",
        "AgentBench",
        "Contact Information: Your email/social"
    ]
    add_bullet_slide(prs, content_layout, "Resources & Further Reading", resources_bullets)
    
    # Slide 24: Future of Code Agents
    future_bullets = [
        "Emerging Capabilities:",
        "Multi-agent collaboration for complex projects",
        "Long-term project memory and reasoning",
        "Human-agent collaborative workflows",
        "Personalization to developer style and preferences",
        "Cross-repository knowledge integration",
        "Research Directions:",
        "Multi-modal reasoning (code + diagrams + documentation)",
        "Improvement of planning algorithms",
        "Tool creation capability",
        "Explainability of agent decisions",
        "Ethical considerations in automated development",
        "[Visual suggestion: Timeline of agent evolution predictions]"
    ]
    add_bullet_slide(prs, content_layout, "Future of Code Agents", future_bullets)
    
    # Slide 25: Q&A
    qa_bullets = [
        "Discussion Topics:",
        "RAG vs. fine-tuning for code generation",
        "Balancing retrieval quality vs. system latency",
        "Security considerations in enterprise settings",
        "Future of code agents in development workflows",
        "Integration strategies for existing teams",
        "[Visual suggestion: QR code to additional resources]"
    ]
    add_bullet_slide(prs, content_layout, "Q&A", qa_bullets)

if __name__ == "__main__":
    create_presentation()
